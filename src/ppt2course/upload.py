"""STEP 1: Upload PPT — parse slide text / speaker notes."""

from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

ROW_TOLERANCE_EMU = Pt(10)


class PptParseError(Exception):
    pass


@dataclass(frozen=True)
class SlideContent:
    index: int
    text: str
    notes: str
    # Excluded from equality so index/text/notes-only callers stay compatible.
    title: str = field(default="", compare=False)
    image_count: int = field(default=0, compare=False)
    has_chart: bool = field(default=False, compare=False)
    shape_count: int = field(default=0, compare=False)


def _shape_text_block(shape) -> str:
    if getattr(shape, "has_table", False):
        rows_text = [
            " ".join(cell.text for cell in row.cells) for row in shape.table.rows
        ]
        return "\n".join(rows_text).strip()

    if getattr(shape, "has_text_frame", False):
        return shape.text_frame.text.strip()

    return ""


def _order_shapes_by_position(shapes) -> list[str]:
    entries = []
    for shape in shapes:
        text = _shape_text_block(shape)
        if text:
            entries.append((shape.top, shape.left, text))

    entries.sort(key=lambda e: (e[0], e[1]))

    rows: list[list[tuple[int, str]]] = []
    row_ref_top = None
    for top, left, text in entries:
        if row_ref_top is None or abs(top - row_ref_top) > ROW_TOLERANCE_EMU:
            rows.append([])
            row_ref_top = top
        rows[-1].append((left, text))

    blocks = []
    for row in rows:
        row.sort(key=lambda e: e[0])
        blocks.extend(text for _, text in row)

    return blocks


def _slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def parse_ppt(path: str) -> list[SlideContent]:
    try:
        presentation = Presentation(path)
    except Exception as exc:
        raise PptParseError(f"failed to open PPT file {path!r}: {exc}") from exc

    result = []
    for index, slide in enumerate(presentation.slides, start=1):
        blocks = _order_shapes_by_position(slide.shapes)
        text = "\n".join(blocks)
        notes = _slide_notes(slide)
        title_shape = slide.shapes.title
        title = _shape_text_block(title_shape) if title_shape is not None else ""
        image_count = sum(
            shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)
            for shape in slide.shapes
        )
        result.append(
            SlideContent(
                index=index,
                text=text,
                notes=notes,
                title=title,
                image_count=image_count,
                has_chart=any(getattr(shape, "has_chart", False) for shape in slide.shapes),
                shape_count=len(slide.shapes),
            )
        )

    return result
