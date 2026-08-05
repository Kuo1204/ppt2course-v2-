"""Integration test against real LibreOffice (STEP 1 preview helper)."""

import pytest
from pptx import Presentation
from pptx.util import Pt

from ppt2course.pptx_preview import PptxPreviewError, find_soffice, render_pptx_thumbnails

HAS_SOFFICE = find_soffice() is not None

pytestmark = pytest.mark.skipif(not HAS_SOFFICE, reason="LibreOffice (soffice) not available")

BLANK_LAYOUT = 6


def _make_pptx(tmp_path, slide_texts, name="deck.pptx"):
    prs = Presentation()
    for text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
        box = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(400), Pt(80))
        box.text_frame.text = text
    path = tmp_path / name
    prs.save(str(path))
    return str(path)


def test_renders_one_png_thumbnail_per_slide_in_order(tmp_path):
    path = _make_pptx(tmp_path, ["第一頁", "第二頁", "第三頁"])

    thumbnails = render_pptx_thumbnails(path)

    assert len(thumbnails) == 3
    for png_bytes in thumbnails:
        assert png_bytes[:8] == b"\x89PNG\r\n\x1a\n"  # real PNG file signature


def test_thumbnail_width_matches_requested_width(tmp_path):
    path = _make_pptx(tmp_path, ["單頁投影片"])

    thumbnails = render_pptx_thumbnails(path, thumbnail_width=320)

    import fitz

    pix = fitz.Pixmap(thumbnails[0])
    assert pix.width == 320


def test_single_slide_deck_returns_one_thumbnail(tmp_path):
    path = _make_pptx(tmp_path, ["只有一頁"])

    thumbnails = render_pptx_thumbnails(path)

    assert len(thumbnails) == 1


def test_raises_pptx_preview_error_for_a_corrupt_file(tmp_path):
    bad = tmp_path / "not_really_a_pptx.pptx"
    # Plain text content is *not* a good "corrupt file" case here — LibreOffice
    # is lenient enough to sniff it as a Writer document and convert it to a
    # one-page PDF regardless of the .pptx extension. Opaque binary garbage
    # is what actually fails to load for it.
    bad.write_bytes(bytes(range(256)) * 20)

    with pytest.raises(PptxPreviewError):
        render_pptx_thumbnails(str(bad))


def test_find_soffice_returns_a_real_executable_path():
    path = find_soffice()
    assert path is not None
    import os

    assert os.path.exists(path)
