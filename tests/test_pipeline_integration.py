"""Real end-to-end integration test: a genuine .pptx through STEP1-6.

Uses NOTES mode so it needs no Gemini key — only real ffmpeg + real edge-tts.
"""

import glob
import os
import shutil
import subprocess

import pytest
from pptx import Presentation
from pptx.util import Pt

from ppt2course.pipeline import run_pipeline
from ppt2course.script_gen import ScriptMode


def _ensure_ffmpeg_on_path() -> bool:
    if shutil.which("ffmpeg") and shutil.which("ffprobe"):
        return True

    candidates = glob.glob(
        os.path.expandvars(
            r"%LOCALAPPDATA%\Microsoft\WinGet\Packages\Gyan.FFmpeg_*\ffmpeg-*\bin"
        )
    )
    for bin_dir in candidates:
        if os.path.exists(os.path.join(bin_dir, "ffmpeg.exe")):
            os.environ["PATH"] = bin_dir + os.pathsep + os.environ.get("PATH", "")
            return True

    return False


HAS_FFMPEG = _ensure_ffmpeg_on_path()

pytestmark = pytest.mark.skipif(not HAS_FFMPEG, reason="ffmpeg/ffprobe not available")


def _make_test_pptx(path: str) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank

    slide1 = prs.slides.add_slide(layout)
    box1 = slide1.shapes.add_textbox(Pt(50), Pt(50), Pt(400), Pt(100))
    box1.text_frame.text = "課程介紹"
    slide1.notes_slide.notes_text_frame.text = "大家好，歡迎來到本次課程。"

    slide2 = prs.slides.add_slide(layout)
    box2 = slide2.shapes.add_textbox(Pt(50), Pt(50), Pt(400), Pt(100))
    box2.text_frame.text = "課程總結"
    slide2.notes_slide.notes_text_frame.text = "今天的課程就到這裡，謝謝大家。"

    prs.save(path)


def _make_color_image(path: str, color: str) -> None:
    subprocess.run(
        ["ffmpeg", "-y", "-f", "lavfi", "-i", f"color=c={color}:s=320x240", "-frames:v", "1", path],
        check=True, capture_output=True,
    )


def test_run_pipeline_end_to_end_with_real_pptx_ffmpeg_and_edge_tts(tmp_path):
    pptx_path = str(tmp_path / "deck.pptx")
    _make_test_pptx(pptx_path)

    img1 = str(tmp_path / "slide1.png")
    img2 = str(tmp_path / "slide2.png")
    _make_color_image(img1, "red")
    _make_color_image(img2, "blue")

    work_dir = str(tmp_path / "work")
    out_dir = str(tmp_path / "out")

    result = run_pipeline(
        pptx_path,
        [img1, img2],
        work_dir,
        out_dir,
        "課程",
        ScriptMode.NOTES,
        "zh-TW-HsiaoChenNeural",
    )

    assert os.path.exists(result["mp4"])
    assert os.path.getsize(result["mp4"]) > 0
    assert os.path.exists(result["srt"])
    assert os.path.exists(result["docx"])

    srt_text = open(result["srt"], encoding="utf-8").read()
    assert "大家好" in srt_text
    assert "謝謝大家" in srt_text

    # intermediate per-slide audio files should be left in work_dir (not cleaned up)
    assert os.path.exists(os.path.join(work_dir, "slide_001.mp3"))
    assert os.path.exists(os.path.join(work_dir, "slide_002.mp3"))
