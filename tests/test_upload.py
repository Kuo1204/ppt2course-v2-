import pytest
from pptx import Presentation
from pptx.util import Pt

from ppt2course.upload import PptParseError, SlideContent, parse_ppt

BLANK_LAYOUT = 6


def new_presentation():
    return Presentation()


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])


def add_text(slide, text, left, top, width=Pt(200), height=Pt(50)):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.text = text
    return box


def add_table(slide, rows_data, left, top, width=Pt(300), height=Pt(100)):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table
    for r, row_values in enumerate(rows_data):
        for c, value in enumerate(row_values):
            table.cell(r, c).text = value
    return graphic_frame


def save(prs, tmp_path, name="test.pptx"):
    path = tmp_path / name
    prs.save(str(path))
    return str(path)


def test_extracts_single_textbox_content(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_text(slide, "你好世界", left=Pt(50), top=Pt(50))
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result == [SlideContent(index=1, text="你好世界", notes="")]


def test_multiple_slides_indexed_from_1(tmp_path):
    prs = new_presentation()
    slide1 = add_slide(prs)
    add_text(slide1, "第一頁", left=Pt(0), top=Pt(0))
    slide2 = add_slide(prs)
    add_text(slide2, "第二頁", left=Pt(0), top=Pt(0))
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert [s.index for s in result] == [1, 2]
    assert result[0].text == "第一頁"
    assert result[1].text == "第二頁"


def test_orders_textboxes_by_row_tolerance_then_left(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_text(slide, "C", left=Pt(20), top=Pt(10))
    add_text(slide, "B", left=Pt(200), top=Pt(15))  # within 10pt of row ref (10) -> same row
    add_text(slide, "A", left=Pt(50), top=Pt(200))  # far below -> new row
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].text == "C\nB\nA"


def test_row_tolerance_inclusive_boundary(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_text(slide, "X", left=Pt(100), top=Pt(0))
    add_text(slide, "Y", left=Pt(10), top=Pt(10))  # diff exactly 10pt -> same row
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].text == "Y\nX"


def test_row_tolerance_just_over_boundary_is_new_row(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_text(slide, "X", left=Pt(100), top=Pt(0))
    add_text(slide, "Y", left=Pt(10), top=Pt(10.1))  # diff > 10pt -> new row
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].text == "X\nY"


def test_extracts_table_cells_space_within_row_newline_between_rows(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_table(slide, [["A1", "B1"], ["A2", "B2"]], left=Pt(0), top=Pt(0))
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].text == "A1 B1\nA2 B2"


def test_textbox_and_table_combined_position_sorted(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_table(slide, [["A1", "B1"]], left=Pt(0), top=Pt(100))
    add_text(slide, "標題", left=Pt(0), top=Pt(0))
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].text == "標題\nA1 B1"


def test_skips_shapes_with_empty_text(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_text(slide, "", left=Pt(0), top=Pt(0))
    add_text(slide, "重點", left=Pt(0), top=Pt(50))
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].text == "重點"


def test_extracts_speaker_notes(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_text(slide, "投影片內容", left=Pt(0), top=Pt(0))
    slide.notes_slide.notes_text_frame.text = "這是備忘稿"
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].notes == "這是備忘稿"


def test_extracts_backward_compatible_visual_metadata(tmp_path):
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
    slide.shapes.title.text = "案例分析"
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].title == "案例分析"
    assert result[0].image_count == 0
    assert result[0].has_chart is False
    assert result[0].shape_count >= 1
    # Metadata does not break legacy equality expectations.
    assert result[0] == SlideContent(index=1, text="案例分析", notes="")


def test_no_notes_slide_returns_empty_string_notes(tmp_path):
    prs = new_presentation()
    slide = add_slide(prs)
    add_text(slide, "投影片內容", left=Pt(0), top=Pt(0))
    path = save(prs, tmp_path)

    result = parse_ppt(path)

    assert result[0].notes == ""


def test_raises_ppt_parse_error_on_missing_file(tmp_path):
    missing = tmp_path / "does_not_exist.pptx"
    with pytest.raises(PptParseError):
        parse_ppt(str(missing))


def test_raises_ppt_parse_error_on_invalid_pptx(tmp_path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"this is not a real pptx file")
    with pytest.raises(PptParseError):
        parse_ppt(str(bad))
