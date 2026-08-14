import base64
import io
import json
from unittest.mock import patch
from urllib.parse import quote

from fastapi.testclient import TestClient

from ppt2course.jobs import JobManager
from ppt2course.media_search import MediaSearchResult
from ppt2course.pipeline import PipelineError
from ppt2course.pptx_preview import PptxPreviewError
from ppt2course.script_gen import ScriptGenerationError, ScriptMode
from ppt2course.server import create_app
from ppt2course.subtitle import TimedChunk
from ppt2course.timeline_models import VisualAsset, VisualAssetType, VisualRecommendation
from ppt2course.tts import TtsError
from ppt2course.upload import PptParseError, SlideContent
from ppt2course.visual_analyzer import DEFAULT_BROLL_DURATION_MS, VisualAnalysisResult


def _make_client(pipeline_fn, tmp_path):
    manager = JobManager(pipeline_fn=pipeline_fn, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    return TestClient(app), manager


def _upload_files():
    return {
        "pptx": ("deck.pptx", io.BytesIO(b"fake-pptx-bytes"), "application/octet-stream"),
        "images": ("slide1.png", io.BytesIO(b"fake-png-bytes"), "image/png"),
    }


def _upload_form(**overrides):
    form = {"script_mode": "NOTES", "voice": "zh-TW-HsiaoChenNeural"}
    form.update(overrides)
    return form


def test_create_job_saves_uploads_and_returns_job_id(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {"mp4": "課程.mp4", "srt": "課程.srt", "docx": "課程.docx"}

    client, manager = _make_client(fake_pipeline, tmp_path)

    response = client.post("/api/jobs", data=_upload_form(), files=_upload_files())

    assert response.status_code == 200
    job_id = response.json()["job_id"]
    assert job_id

    manager.process_next()

    assert len(calls) == 1
    kwargs = calls[0]
    assert kwargs["pptx_path"].endswith("deck.pptx")
    import os
    assert os.path.exists(kwargs["pptx_path"])
    assert len(kwargs["image_paths"]) == 1
    assert os.path.exists(kwargs["image_paths"][0])


def test_create_job_forwards_voice_rate_and_volume(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)

    client.post(
        "/api/jobs",
        data=_upload_form(rate="+20%", volume="-10%", font_size="60"),
        files=_upload_files(),
    )
    manager.process_next()

    assert calls[0]["voice_rate"] == "+20%"
    assert calls[0]["voice_volume"] == "-10%"
    assert calls[0]["font_size"] == 60


def test_create_job_forwards_subtitle_margin_v(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)

    client.post(
        "/api/jobs",
        data=_upload_form(subtitle_margin_v="220"),
        files=_upload_files(),
    )
    manager.process_next()

    assert calls[0]["subtitle_margin_v"] == 220


def test_create_job_forwards_logo_position(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)

    client.post(
        "/api/jobs",
        data=_upload_form(logo_position="bottom-left"),
        files=_upload_files(),
    )
    manager.process_next()

    assert calls[0]["logo_position"] == "bottom-left"


def test_create_job_defaults_logo_position_to_top_right(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)

    client.post("/api/jobs", data=_upload_form(), files=_upload_files())
    manager.process_next()

    assert calls[0]["logo_position"] == "top-right"


def test_create_job_rejects_invalid_logo_position(tmp_path):
    def fake_pipeline(**kwargs):
        return {}

    client, _manager = _make_client(fake_pipeline, tmp_path)

    response = client.post(
        "/api/jobs",
        data=_upload_form(logo_position="middle"),
        files=_upload_files(),
    )

    assert response.status_code == 400


def test_create_job_writes_custom_dict_to_a_file_and_forwards_its_path(tmp_path):
    import os

    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)

    client.post(
        "/api/jobs",
        data=_upload_form(custom_dict="普拉斯提亞 100\n辰昕科技\n"),
        files=_upload_files(),
    )
    manager.process_next()

    custom_dict_path = calls[0]["custom_dict_path"]
    assert custom_dict_path is not None
    assert os.path.exists(custom_dict_path)
    assert "普拉斯提亞 100" in open(custom_dict_path, encoding="utf-8").read()


def test_create_job_leaves_custom_dict_path_none_when_not_supplied(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)

    client.post("/api/jobs", data=_upload_form(), files=_upload_files())
    manager.process_next()

    assert calls[0]["custom_dict_path"] is None


def test_create_job_forwards_logo_opacity(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)

    client.post(
        "/api/jobs",
        data=_upload_form(logo_opacity="0.35"),
        files={**_upload_files(), "logo": ("logo.png", io.BytesIO(b"fake-png"), "image/png")},
    )
    manager.process_next()

    assert calls[0]["logo_opacity"] == 0.35
    assert calls[0]["logo_path"] is not None


def test_create_job_defaults_logo_opacity_to_fully_opaque(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)

    client.post("/api/jobs", data=_upload_form(), files=_upload_files())
    manager.process_next()

    assert calls[0]["logo_opacity"] == 1.0


# ---- broll_selections (confirmed B-roll picks, downloaded server-side) ----


def test_create_job_defaults_broll_selections_to_empty_list(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    client.post("/api/jobs", data=_upload_form(), files=_upload_files())
    manager.process_next()

    assert calls[0]["broll_selections"] == []


def test_create_job_downloads_and_forwards_broll_selection(tmp_path):
    import os

    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    selections = json.dumps(
        [{"slide_number": 1, "download_url": "https://example.com/a.jpg", "start_ms": 100, "end_ms": 500}]
    )
    with patch("ppt2course.server.urlopen") as mock_urlopen:
        mock_urlopen.return_value.__enter__.return_value.read.return_value = b"fake-jpeg-bytes"
        client.post(
            "/api/jobs",
            data=_upload_form(broll_selections=selections),
            files=_upload_files(),
        )
    manager.process_next()

    resolved = calls[0]["broll_selections"]
    assert len(resolved) == 1
    assert resolved[0]["slide_number"] == 1
    assert resolved[0]["start_ms"] == 100
    assert resolved[0]["end_ms"] == 500
    assert os.path.exists(resolved[0]["image_path"])
    assert open(resolved[0]["image_path"], "rb").read() == b"fake-jpeg-bytes"


def test_create_job_skips_broll_selection_when_download_fails(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    selections = json.dumps(
        [{"slide_number": 1, "download_url": "https://example.com/a.jpg", "start_ms": 100, "end_ms": 500}]
    )
    with patch("ppt2course.server.urlopen", side_effect=OSError("network error")):
        response = client.post(
            "/api/jobs",
            data=_upload_form(broll_selections=selections),
            files=_upload_files(),
        )
    manager.process_next()

    assert response.status_code == 200  # job creation itself never fails
    assert calls[0]["broll_selections"] == []


def test_create_job_skips_malformed_broll_selection(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    response = client.post(
        "/api/jobs",
        data=_upload_form(broll_selections="not valid json"),
        files=_upload_files(),
    )
    manager.process_next()

    assert response.status_code == 200
    assert calls[0]["broll_selections"] == []


# ---- avatar_mode / avatar_position / avatar_size / avatar_custom_slides ----


def test_create_job_defaults_avatar_mode_to_none(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    client.post("/api/jobs", data=_upload_form(), files=_upload_files())
    manager.process_next()

    assert calls[0]["avatar_mode"] == "none"
    assert calls[0]["avatar_position"] == "bottom_right"
    assert calls[0]["avatar_size"] == "small"
    assert calls[0]["avatar_custom_slides"] == []


def test_create_job_forwards_avatar_settings(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    client.post(
        "/api/jobs",
        data=_upload_form(
            avatar_mode="custom",
            avatar_position="left",
            avatar_size="large",
            avatar_margin="10",
            avatar_custom_slides="[1, 3]",
        ),
        files=_upload_files(),
    )
    manager.process_next()

    assert calls[0]["avatar_mode"] == "custom"
    assert calls[0]["avatar_position"] == "left"
    assert calls[0]["avatar_size"] == "large"
    assert calls[0]["avatar_margin"] == 10
    assert calls[0]["avatar_custom_slides"] == [1, 3]


def test_create_job_rejects_invalid_avatar_mode(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs", data=_upload_form(avatar_mode="bogus"), files=_upload_files()
    )
    assert response.status_code == 400


def test_create_job_rejects_invalid_avatar_position(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs", data=_upload_form(avatar_position="top-center"), files=_upload_files()
    )
    assert response.status_code == 400


def test_create_job_rejects_invalid_avatar_size(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs", data=_upload_form(avatar_size="huge"), files=_upload_files()
    )
    assert response.status_code == 400


def test_create_job_rejects_malformed_avatar_custom_slides(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs",
        data=_upload_form(avatar_mode="custom", avatar_custom_slides="not json"),
        files=_upload_files(),
    )
    assert response.status_code == 400


def test_create_job_rejects_non_integer_avatar_custom_slides(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs",
        data=_upload_form(avatar_mode="custom", avatar_custom_slides='["a", "b"]'),
        files=_upload_files(),
    )
    assert response.status_code == 400


# ---- reading_pause_ms / closing_pause_ms / target_duration_ms / enable_ken_burns ----


def test_create_job_defaults_pacing_settings(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    client.post("/api/jobs", data=_upload_form(), files=_upload_files())
    manager.process_next()

    assert calls[0]["reading_pause_ms"] == 0
    assert calls[0]["closing_pause_ms"] == 0
    assert calls[0]["target_duration_ms"] is None
    assert calls[0]["enable_ken_burns"] is False


def test_create_job_forwards_pacing_settings(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    client.post(
        "/api/jobs",
        data=_upload_form(
            reading_pause_ms="800",
            closing_pause_ms="2000",
            target_duration_ms="60000",
            enable_ken_burns="true",
        ),
        files=_upload_files(),
    )
    manager.process_next()

    assert calls[0]["reading_pause_ms"] == 800
    assert calls[0]["closing_pause_ms"] == 2000
    assert calls[0]["target_duration_ms"] == 60000
    assert calls[0]["enable_ken_burns"] is True


def test_create_job_rejects_negative_reading_pause_ms(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs", data=_upload_form(reading_pause_ms="-1"), files=_upload_files()
    )
    assert response.status_code == 400


def test_create_job_rejects_negative_closing_pause_ms(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs", data=_upload_form(closing_pause_ms="-1"), files=_upload_files()
    )
    assert response.status_code == 400


def test_create_job_rejects_non_positive_target_duration_ms(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs", data=_upload_form(target_duration_ms="0"), files=_upload_files()
    )
    assert response.status_code == 400


def test_get_job_status_surfaces_target_duration_reachable(tmp_path):
    client, manager = _make_client(
        lambda **kwargs: {"target_duration_reachable": False}, tmp_path
    )
    response = client.post(
        "/api/jobs", data=_upload_form(target_duration_ms="1000"), files=_upload_files()
    )
    job_id = response.json()["job_id"]
    manager.process_next()

    body = client.get(f"/api/jobs/{job_id}").json()
    assert body["details"]["target_duration_reachable"] is False


def test_get_job_status_target_duration_reachable_absent_when_not_requested(tmp_path):
    client, manager = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post("/api/jobs", data=_upload_form(), files=_upload_files())
    job_id = response.json()["job_id"]
    manager.process_next()

    body = client.get(f"/api/jobs/{job_id}").json()
    assert body["details"]["target_duration_reachable"] is None


def test_get_job_status_queued_before_processing(tmp_path):
    client, manager = _make_client(lambda **kwargs: {}, tmp_path)

    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]

    response = client.get(f"/api/jobs/{job_id}")
    assert response.status_code == 200
    assert response.json() == {"status": "queued"}


def test_get_job_status_done_includes_download_links(tmp_path):
    def fake_pipeline(**kwargs):
        return {"mp4": "課程.mp4", "srt": "課程.srt", "docx": "課程.docx"}

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}")
    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "done"
    assert body["downloads"] == {
        "mp4": f"/api/jobs/{job_id}/download/mp4",
        "srt": f"/api/jobs/{job_id}/download/srt",
        "docx": f"/api/jobs/{job_id}/download/docx",
    }


def test_get_job_status_done_includes_video_and_script_details(tmp_path):
    def fake_pipeline(**kwargs):
        return {
            "mp4": "課程.mp4",
            "srt": "課程.srt",
            "docx": "課程.docx",
            "video_size_bytes": 5_242_880,
            "video_duration_ms": 65_000,
            "script_char_count": 812,
        }

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}")
    assert response.status_code == 200
    body = response.json()
    assert body["details"]["video_size_bytes"] == 5_242_880
    assert body["details"]["video_duration_ms"] == 65_000
    assert body["details"]["script_char_count"] == 812
    # generation_seconds is derived from Job.started_at/completed_at, which
    # process_next() sets to real wall-clock times around the (instant, in
    # this test) fake pipeline call — must be a small non-negative number,
    # not None and not the actual seconds since epoch.
    assert isinstance(body["details"]["generation_seconds"], (int, float))
    assert 0 <= body["details"]["generation_seconds"] < 5


def test_get_job_status_error_includes_message(tmp_path):
    def failing_pipeline(**kwargs):
        raise PipelineError("TTS synthesis failed for slide 1: network error")

    client, manager = _make_client(failing_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}")
    assert response.status_code == 200
    assert response.json() == {
        "status": "error",
        "error": "TTS synthesis failed for slide 1: network error",
    }


def test_get_unknown_job_returns_404(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.get("/api/jobs/does-not-exist")
    assert response.status_code == 404


def test_invalid_script_mode_returns_400(tmp_path):
    client, _ = _make_client(lambda **kwargs: {}, tmp_path)
    response = client.post(
        "/api/jobs", data=_upload_form(script_mode="NOT_A_MODE"), files=_upload_files()
    )
    assert response.status_code == 400


def test_download_file_after_done(tmp_path):
    out_file = tmp_path / "課程.mp4"
    out_file.write_bytes(b"video-bytes")

    def fake_pipeline(**kwargs):
        return {"mp4": str(out_file)}

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}/download/mp4")
    assert response.status_code == 200
    assert response.content == b"video-bytes"


def test_download_file_uses_the_course_name_not_a_generic_browser_filename(tmp_path):
    # FileResponse without an explicit filename= has no Content-Disposition
    # filename, so the browser falls back to a generic name like "mp4 (1).mp4"
    # derived from the URL — not the user's chosen course name.
    out_file = tmp_path / "試作影片.mp4"
    out_file.write_bytes(b"video-bytes")

    def fake_pipeline(**kwargs):
        return {"mp4": str(out_file)}

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}/download/mp4")
    assert response.status_code == 200
    content_disposition = response.headers["content-disposition"]
    # non-ASCII filenames are RFC 5987 percent-encoded in the header
    assert quote("試作影片.mp4") in content_disposition


def test_download_before_done_returns_404(tmp_path):
    client, manager = _make_client(lambda **kwargs: {"mp4": "x"}, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]

    response = client.get(f"/api/jobs/{job_id}/download/mp4")
    assert response.status_code == 404


def test_download_unknown_filetype_returns_404(tmp_path):
    def fake_pipeline(**kwargs):
        return {"mp4": "x", "srt": "y", "docx": "z"}

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}/download/exe")
    assert response.status_code == 404


# ---------- inline "view" (video preview) and the /share landing page ----------
# Scanning the QR code used to link straight to the attachment-disposition
# download route, so the phone's browser silently triggered a download and
# left a blank tab behind — no confirmation, no way to preview first. /share
# gives the phone something to actually land on: a page that asks whether to
# preview or download.


def test_view_file_after_done_uses_inline_not_attachment_disposition(tmp_path):
    out_file = tmp_path / "課程.mp4"
    out_file.write_bytes(b"video-bytes")

    def fake_pipeline(**kwargs):
        return {"mp4": str(out_file)}

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}/view/mp4")
    assert response.status_code == 200
    assert response.content == b"video-bytes"
    disposition = response.headers.get("content-disposition", "")
    assert not disposition.startswith("attachment")


def test_view_before_done_returns_404(tmp_path):
    client, manager = _make_client(lambda **kwargs: {"mp4": "x"}, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]

    response = client.get(f"/api/jobs/{job_id}/view/mp4")
    assert response.status_code == 404


def test_view_unknown_filetype_returns_404(tmp_path):
    def fake_pipeline(**kwargs):
        return {"mp4": "x"}

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}/view/exe")
    assert response.status_code == 404


def test_share_page_when_done_offers_preview_and_download(tmp_path):
    out_file = tmp_path / "課程.mp4"
    out_file.write_bytes(b"video-bytes")

    def fake_pipeline(**kwargs):
        return {"mp4": str(out_file)}

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/share/{job_id}")

    assert response.status_code == 200
    assert response.headers["content-type"].startswith("text/html")
    assert f"/api/jobs/{job_id}/view/mp4" in response.text
    assert f"/api/jobs/{job_id}/download/mp4" in response.text


def test_share_page_still_processing_shows_wait_message_not_a_blank_page(tmp_path):
    client, manager = _make_client(lambda **kwargs: {"mp4": "x"}, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]

    response = client.get(f"/share/{job_id}")

    assert response.status_code == 200
    assert "製作中" in response.text


def test_share_page_error_status_shows_the_error_message(tmp_path):
    def fake_pipeline(**kwargs):
        raise PipelineError("boom")

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post("/api/jobs", data=_upload_form(), files=_upload_files()).json()["job_id"]
    manager.process_next()

    response = client.get(f"/share/{job_id}")

    assert response.status_code == 200
    assert "boom" in response.text


def test_share_page_unknown_job_returns_a_friendly_404_page(tmp_path):
    client, manager = _make_client(lambda **kwargs: {}, tmp_path)

    response = client.get("/share/does-not-exist")

    assert response.status_code == 404
    assert response.headers["content-type"].startswith("text/html")


def test_view_and_share_routes_are_exempt_from_basic_auth(tmp_path):
    out_file = tmp_path / "課程.mp4"
    out_file.write_bytes(b"video-bytes")

    def fake_pipeline(**kwargs):
        return {"mp4": str(out_file)}

    manager = JobManager(pipeline_fn=fake_pipeline, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        basic_auth_user="admin",
        basic_auth_password="s3cret",
    )
    client = TestClient(app)
    job_id = client.post(
        "/api/jobs", data=_upload_form(), files=_upload_files(),
        headers=_basic_auth_header("admin", "s3cret"),
    ).json()["job_id"]
    manager.process_next()

    assert client.get(f"/api/jobs/{job_id}/view/mp4").status_code == 200
    assert client.get(f"/share/{job_id}").status_code == 200


def test_serves_built_frontend_index_at_root(tmp_path):
    dist_dir = tmp_path / "dist"
    dist_dir.mkdir()
    (dist_dir / "index.html").write_text("<html>PPT2COURSE_FRONTEND</html>", encoding="utf-8")

    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(
        job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=str(dist_dir)
    )
    client = TestClient(app)

    response = client.get("/")
    assert response.status_code == 200
    assert "PPT2COURSE_FRONTEND" in response.text


def test_serves_built_frontend_static_assets(tmp_path):
    dist_dir = tmp_path / "dist"
    assets_dir = dist_dir / "assets"
    assets_dir.mkdir(parents=True)
    (dist_dir / "index.html").write_text("<html></html>", encoding="utf-8")
    (assets_dir / "app.js").write_text("console.log('hi')", encoding="utf-8")

    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(
        job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=str(dist_dir)
    )
    client = TestClient(app)

    response = client.get("/assets/app.js")
    assert response.status_code == 200
    assert "console.log" in response.text


def test_api_routes_still_work_when_frontend_is_mounted(tmp_path):
    dist_dir = tmp_path / "dist"
    dist_dir.mkdir()
    (dist_dir / "index.html").write_text("<html></html>", encoding="utf-8")

    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(
        job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=str(dist_dir)
    )
    client = TestClient(app)

    response = client.get("/api/jobs/does-not-exist")
    assert response.status_code == 404
    assert response.json() == {"detail": "job not found"}


def test_missing_frontend_dist_does_not_crash_app_creation(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=str(tmp_path / "no-such-dir"),
    )
    client = TestClient(app)

    response = client.get("/")
    assert response.status_code == 404


def test_extract_script_text_from_txt_upload(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    response = client.post(
        "/api/extract-script-text",
        files={"file": ("script.txt", io.BytesIO("第1頁\n大家好".encode("utf-8")), "text/plain")},
    )
    assert response.status_code == 200
    assert response.json() == {"text": "第1頁\n大家好"}


def test_extract_script_text_rejects_unsupported_extension(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    response = client.post(
        "/api/extract-script-text",
        files={"file": ("script.pdf", io.BytesIO(b"whatever"), "application/pdf")},
    )
    assert response.status_code == 400


# ---------- AI script preview (AUTO / POLISH) ----------
# Generating the whole video just to find out what the AI wrote meant the
# user couldn't review or fix a bad script until the entire pipeline had
# already run. This endpoint runs only the Gemini call, synchronously, so
# the UI can show the result and let the user continue (or not) before ever
# submitting a job.


def test_generate_script_preview_auto_mode_returns_generated_texts(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch(
            "ppt2course.server.generate_script", return_value=["AI 寫的講稿"]
        ) as mock_generate:
            response = client.post(
                "/api/generate-script",
                data={"script_mode": "AUTO", "gemini_api_key": "secret-key"},
                files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
            )

    assert response.status_code == 200
    assert response.json() == {"texts": ["AI 寫的講稿"]}
    assert mock_generate.call_args.args[0] == ScriptMode.AUTO
    assert mock_generate.call_args.kwargs["api_key"] == "secret-key"


def test_generate_script_preview_polish_mode_forwards_input_texts(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch(
            "ppt2course.server.generate_script", return_value=["潤飾後的講稿"]
        ) as mock_generate:
            response = client.post(
                "/api/generate-script",
                data={
                    "script_mode": "POLISH",
                    "gemini_api_key": "secret-key",
                    "texts": json.dumps(["我自己寫的草稿"]),
                },
                files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
            )

    assert response.status_code == 200
    assert response.json() == {"texts": ["潤飾後的講稿"]}
    assert mock_generate.call_args.kwargs["texts"] == ["我自己寫的草稿"]


def test_generate_script_preview_rejects_notes_and_own_modes(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    for mode in ("NOTES", "OWN"):
        response = client.post(
            "/api/generate-script",
            data={"script_mode": mode, "gemini_api_key": "secret-key"},
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )
        assert response.status_code == 400


def test_generate_script_preview_invalid_script_mode_returns_400(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    response = client.post(
        "/api/generate-script",
        data={"script_mode": "NOT_A_MODE", "gemini_api_key": "secret-key"},
        files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
    )
    assert response.status_code == 400


def test_generate_script_preview_wraps_ppt_parse_error(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    with patch("ppt2course.server.parse_ppt", side_effect=PptParseError("bad pptx")):
        response = client.post(
            "/api/generate-script",
            data={"script_mode": "AUTO", "gemini_api_key": "secret-key"},
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )
    assert response.status_code == 400


def test_generate_script_preview_wraps_gemini_error(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch(
            "ppt2course.server.generate_script",
            side_effect=ScriptGenerationError("Gemini API call failed: boom"),
        ):
            response = client.post(
                "/api/generate-script",
                data={"script_mode": "AUTO", "gemini_api_key": "secret-key"},
                files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
            )
    assert response.status_code == 502
    assert "boom" in response.json()["detail"]


def test_generate_script_preview_never_stores_the_api_key(tmp_path):
    # Same guarantee job creation already gives: the key reaches the call
    # but is never echoed back in the response.
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch("ppt2course.server.generate_script", return_value=["text"]):
            response = client.post(
                "/api/generate-script",
                data={"script_mode": "AUTO", "gemini_api_key": "super-secret-key"},
                files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
            )
    assert "super-secret-key" not in response.text


# ---------- AI visual-need analysis (preview-only, never touches compose_video) ----------
# Mirrors /api/generate-script's contract: runs synchronously, nothing is
# persisted, and a Gemini/parse failure surfaces as an HTTP error here rather
# than ever being allowed to fail a whole video job.


def test_analyze_visuals_returns_recommendations(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    fake_result = VisualAnalysisResult(
        recommendations=(
            VisualRecommendation(
                slide_number=1,
                title="職場孤立",
                visual_need_score=82,
                recommended=True,
                reason="抽象概念適合搭配情境圖片",
                visual_type=VisualAssetType.IMAGE,
                keywords=("employee isolation",),
                suggested_position="during_slide",
            ),
        ),
        used_ai=True,
    )
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch(
            "ppt2course.server.analyze_visual_needs", return_value=fake_result
        ) as mock_analyze:
            response = client.post(
                "/api/analyze-visuals",
                data={
                    "texts": json.dumps(["講稿一"]),
                    "voice": "zh-TW-HsiaoChenNeural",
                    "gemini_api_key": "secret-key",
                },
                files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
            )

    assert response.status_code == 200
    body = response.json()
    assert body["used_ai"] is True
    assert body["warnings"] == []
    assert body["recommendations"] == [
        {
            "slide_number": 1,
            "title": "職場孤立",
            "visual_need_score": 82,
            "recommended": True,
            "reason": "抽象概念適合搭配情境圖片",
            "visual_type": "image",
            "keywords": ["employee isolation"],
            "suggested_position": "during_slide",
            "script_anchor": "",
            # No script_anchor -> the fixed default window, no TTS call needed.
            "suggested_start_ms": 0,
            "suggested_end_ms": DEFAULT_BROLL_DURATION_MS,
        }
    ]
    assert mock_analyze.call_args.kwargs["api_key"] == "secret-key"


def test_analyze_visuals_rejects_mismatched_texts_length(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [
        SlideContent(index=1, text="投影片一", notes=""),
        SlideContent(index=2, text="投影片二", notes=""),
    ]
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        response = client.post(
            "/api/analyze-visuals",
            data={"texts": json.dumps(["只有一段講稿"]), "voice": "zh-TW-HsiaoChenNeural"},
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )

    assert response.status_code == 400


def test_analyze_visuals_wraps_ppt_parse_error(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    with patch("ppt2course.server.parse_ppt", side_effect=PptParseError("bad pptx")):
        response = client.post(
            "/api/analyze-visuals",
            data={"texts": json.dumps([]), "voice": "zh-TW-HsiaoChenNeural"},
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )

    assert response.status_code == 400
    assert "bad pptx" in response.json()["detail"]


def test_analyze_visuals_never_echoes_the_api_key(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    fake_result = VisualAnalysisResult(recommendations=(), warnings=("Gemini API Key 未設定，已使用本機規則產生建議。",))
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch("ppt2course.server.analyze_visual_needs", return_value=fake_result):
            response = client.post(
                "/api/analyze-visuals",
                data={
                    "texts": json.dumps(["講稿一"]),
                    "voice": "zh-TW-HsiaoChenNeural",
                    "gemini_api_key": "super-secret-key",
                },
                files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
            )

    assert response.status_code == 200
    assert "super-secret-key" not in response.text
    assert response.json()["warnings"] == ["Gemini API Key 未設定，已使用本機規則產生建議。"]


# ---- suggested_start_ms/end_ms (real per-slide TTS timing, opt-in per recommendation) ----


def _fake_recommendation(recommended=True, script_anchor="孤立"):
    return VisualRecommendation(
        slide_number=1,
        title="職場孤立",
        visual_need_score=82 if recommended else 20,
        recommended=recommended,
        reason="抽象概念適合搭配情境圖片",
        visual_type=VisualAssetType.IMAGE,
        keywords=("employee isolation",),
        suggested_position="during_slide",
        script_anchor=script_anchor,
    )


def test_analyze_visuals_uses_real_tts_timing_for_a_recommended_slide_with_anchor(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    fake_result = VisualAnalysisResult(recommendations=(_fake_recommendation(),))
    fake_chunks = [
        TimedChunk("員工被", 0, 600),
        TimedChunk("孤立", 600, 1100),
        TimedChunk("的情況", 1100, 1800),
    ]
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch("ppt2course.server.analyze_visual_needs", return_value=fake_result):
            with patch(
                "ppt2course.server.synthesize", return_value=fake_chunks
            ) as mock_synthesize:
                response = client.post(
                    "/api/analyze-visuals",
                    data={
                        "texts": json.dumps(["員工被孤立的情況"]),
                        "voice": "zh-TW-HsiaoChenNeural",
                        "voice_rate": "+10%",
                        "voice_volume": "-5%",
                    },
                    files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
                )

    assert response.status_code == 200
    rec = response.json()["recommendations"][0]
    assert rec["suggested_start_ms"] == 600
    assert rec["suggested_end_ms"] == 1100
    assert mock_synthesize.call_args.args[0] == "員工被孤立的情況"
    assert mock_synthesize.call_args.args[1] == "zh-TW-HsiaoChenNeural"
    assert mock_synthesize.call_args.kwargs["rate"] == "+10%"
    assert mock_synthesize.call_args.kwargs["volume"] == "-5%"


def test_analyze_visuals_skips_tts_for_a_non_recommended_slide(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    fake_result = VisualAnalysisResult(
        recommendations=(_fake_recommendation(recommended=False),)
    )
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch("ppt2course.server.analyze_visual_needs", return_value=fake_result):
            with patch("ppt2course.server.synthesize") as mock_synthesize:
                response = client.post(
                    "/api/analyze-visuals",
                    data={"texts": json.dumps(["講稿一"]), "voice": "zh-TW-HsiaoChenNeural"},
                    files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
                )

    mock_synthesize.assert_not_called()
    rec = response.json()["recommendations"][0]
    assert (rec["suggested_start_ms"], rec["suggested_end_ms"]) == (0, DEFAULT_BROLL_DURATION_MS)


def test_analyze_visuals_skips_tts_when_there_is_no_script_anchor(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    fake_result = VisualAnalysisResult(recommendations=(_fake_recommendation(script_anchor=""),))
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch("ppt2course.server.analyze_visual_needs", return_value=fake_result):
            with patch("ppt2course.server.synthesize") as mock_synthesize:
                response = client.post(
                    "/api/analyze-visuals",
                    data={"texts": json.dumps(["講稿一"]), "voice": "zh-TW-HsiaoChenNeural"},
                    files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
                )

    mock_synthesize.assert_not_called()
    rec = response.json()["recommendations"][0]
    assert (rec["suggested_start_ms"], rec["suggested_end_ms"]) == (0, DEFAULT_BROLL_DURATION_MS)


def test_analyze_visuals_falls_back_to_default_window_when_tts_fails(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="")]
    fake_result = VisualAnalysisResult(recommendations=(_fake_recommendation(),))
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        with patch("ppt2course.server.analyze_visual_needs", return_value=fake_result):
            with patch(
                "ppt2course.server.synthesize", side_effect=TtsError("network error")
            ):
                response = client.post(
                    "/api/analyze-visuals",
                    data={"texts": json.dumps(["員工被孤立的情況"]), "voice": "zh-TW-HsiaoChenNeural"},
                    files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
                )

    assert response.status_code == 200  # a TTS failure here must never fail the whole request
    rec = response.json()["recommendations"][0]
    assert (rec["suggested_start_ms"], rec["suggested_end_ms"]) == (0, DEFAULT_BROLL_DURATION_MS)


# ---------- Pexels media search (optional, never fails the video job) ----------


def test_media_search_returns_assets(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_result = MediaSearchResult(
        assets=(
            VisualAsset(
                source="pexels",
                asset_type=VisualAssetType.IMAGE,
                preview_url="https://example.com/preview.jpg",
                download_url="https://example.com/full.jpg",
                keyword="employee isolation",
                slide_number=1,
                photographer="Jane Doe",
            ),
        )
    )
    with patch("ppt2course.server.search_pexels", return_value=fake_result) as mock_search:
        response = client.get(
            "/api/media-search",
            params={"keyword": "employee isolation", "slide_number": 1},
        )

    assert response.status_code == 200
    body = response.json()
    assert body["warning"] is None
    assert body["assets"] == [
        {
            "source": "pexels",
            "asset_type": "image",
            "preview_url": "https://example.com/preview.jpg",
            "download_url": "https://example.com/full.jpg",
            "keyword": "employee isolation",
            "slide_number": 1,
            "photographer": "Jane Doe",
            "local_path": None,
        }
    ]
    assert mock_search.call_args.args[0] == "employee isolation"
    assert mock_search.call_args.args[1] == 1


def test_media_search_returns_warning_without_error_status(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    fake_result = MediaSearchResult(warning="Pexels API Key 未設定，無法搜尋免費素材。")
    with patch("ppt2course.server.search_pexels", return_value=fake_result):
        response = client.get(
            "/api/media-search",
            params={"keyword": "employee isolation", "slide_number": 1},
        )

    assert response.status_code == 200
    assert response.json() == {"assets": [], "warning": "Pexels API Key 未設定，無法搜尋免費素材。"}


def test_media_search_rejects_invalid_media_type(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    response = client.get(
        "/api/media-search",
        params={"keyword": "x", "slide_number": 1, "media_type": "gif"},
    )

    assert response.status_code == 400


def test_media_search_forwards_optional_api_key_and_limit(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    with patch(
        "ppt2course.server.search_pexels", return_value=MediaSearchResult()
    ) as mock_search:
        client.get(
            "/api/media-search",
            params={
                "keyword": "office",
                "slide_number": 3,
                "media_type": "video",
                "pexels_api_key": "my-key",
                "limit": 4,
            },
        )

    assert mock_search.call_args.kwargs["api_key"] == "my-key"
    assert mock_search.call_args.kwargs["limit"] == 4
    assert mock_search.call_args.kwargs["media_type"] is VisualAssetType.VIDEO


# ---------- PPTX preview (real slide thumbnails for the upload step) ----------
# Rendering the whole video (or even just extracting text) doesn't tell the
# user whether they uploaded the right deck. This endpoint runs LibreOffice
# synchronously to produce real slide thumbnails, mirroring the visual
# feedback the per-slide image upload field already gives.


def test_pptx_preview_returns_one_data_url_per_thumbnail():
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=None, frontend_dist=None)
    client = TestClient(app)

    fake_pngs = [b"\x89PNG-page-1", b"\x89PNG-page-2"]
    with patch("ppt2course.server.render_pptx_thumbnails", return_value=fake_pngs) as mock_render:
        response = client.post(
            "/api/pptx-preview",
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )

    assert response.status_code == 200
    body = response.json()
    assert len(body["thumbnails"]) == 2
    for url in body["thumbnails"]:
        assert url.startswith("data:image/png;base64,")
    mock_render.assert_called_once()


def test_pptx_preview_wraps_conversion_error_as_502():
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=None, frontend_dist=None)
    client = TestClient(app)

    with patch(
        "ppt2course.server.render_pptx_thumbnails",
        side_effect=PptxPreviewError("LibreOffice not found"),
    ):
        response = client.post(
            "/api/pptx-preview",
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )

    assert response.status_code == 502
    assert "LibreOffice not found" in response.json()["detail"]


def test_pptx_preview_never_persists_the_uploaded_file(tmp_path):
    data_root = tmp_path / "data"
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(data_root), frontend_dist=None)
    client = TestClient(app)

    with patch("ppt2course.server.render_pptx_thumbnails", return_value=[b"\x89PNG"]):
        client.post(
            "/api/pptx-preview",
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )

    leftover_pptx = list(data_root.rglob("*.pptx")) if data_root.exists() else []
    assert leftover_pptx == []


# ---------- PPTX notes preview (speaker notes for the NOTES script mode) ----------
# NOTES mode uses each slide's speaker-notes text verbatim as the narration
# script, but the user has no way to see that text without reopening the
# deck in PowerPoint. This endpoint parses the already-uploaded .pptx and
# hands back each slide's notes so the frontend can show it before submit.


def test_pptx_notes_returns_notes_text_per_slide():
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=None, frontend_dist=None)
    client = TestClient(app)

    fake_slides = [
        SlideContent(index=1, text="投影片一", notes="第一頁備忘稿"),
        SlideContent(index=2, text="投影片二", notes=""),
        SlideContent(index=3, text="投影片三", notes="第三頁備忘稿"),
    ]
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        response = client.post(
            "/api/pptx-notes",
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )

    assert response.status_code == 200
    assert response.json() == {"notes": ["第一頁備忘稿", "", "第三頁備忘稿"]}


def test_pptx_notes_wraps_parse_error_as_400():
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=None, frontend_dist=None)
    client = TestClient(app)

    with patch("ppt2course.server.parse_ppt", side_effect=PptParseError("corrupt file")):
        response = client.post(
            "/api/pptx-notes",
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )

    assert response.status_code == 400
    assert "corrupt file" in response.json()["detail"]


def test_pptx_notes_extracts_real_speaker_notes_end_to_end(tmp_path):
    # No mocking of parse_ppt here: builds an actual .pptx with python-pptx,
    # posts the real bytes through the HTTP layer, and confirms the notes
    # that come back are exactly what was typed into the speaker-notes box —
    # the same real-tool guarantee the rest of this project holds itself to.
    from pptx import Presentation

    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.notes_slide.notes_text_frame.text = "第一頁的備忘稿內容"
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # no notes added

    pptx_path = tmp_path / "real_deck.pptx"
    prs.save(str(pptx_path))

    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=None, frontend_dist=None)
    client = TestClient(app)

    with open(pptx_path, "rb") as f:
        response = client.post(
            "/api/pptx-notes",
            files={"pptx": ("real_deck.pptx", f, "application/octet-stream")},
        )

    assert response.status_code == 200
    assert response.json() == {"notes": ["第一頁的備忘稿內容", ""]}


def test_pptx_notes_never_persists_the_uploaded_file(tmp_path):
    data_root = tmp_path / "data"
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(data_root), frontend_dist=None)
    client = TestClient(app)

    fake_slides = [SlideContent(index=1, text="投影片一", notes="備忘稿")]
    with patch("ppt2course.server.parse_ppt", return_value=fake_slides):
        client.post(
            "/api/pptx-notes",
            files={"pptx": ("deck.pptx", io.BytesIO(b"fake-pptx"), "application/octet-stream")},
        )

    leftover_pptx = list(data_root.rglob("*.pptx")) if data_root.exists() else []
    assert leftover_pptx == []


def test_voice_preview_returns_audio_and_caches_by_voice(tmp_path):
    calls = []

    def fake_preview(text, voice, rate=None, volume=None):
        calls.append((text, voice, rate, volume))
        return b"fake-mp3-bytes"

    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        voice_preview_fn=fake_preview,
    )
    client = TestClient(app)

    first = client.get("/api/voice-preview/zh-TW-HsiaoChenNeural")
    assert first.status_code == 200
    assert first.content == b"fake-mp3-bytes"
    assert first.headers["content-type"] == "audio/mpeg"

    second = client.get("/api/voice-preview/zh-TW-HsiaoChenNeural")
    assert second.content == b"fake-mp3-bytes"

    assert len(calls) == 1  # cached on the second request, no repeat synthesis


def test_voice_preview_rejects_unknown_voice(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        voice_preview_fn=lambda text, voice, rate=None, volume=None: b"",
    )
    client = TestClient(app)

    response = client.get("/api/voice-preview/not-a-real-voice")
    assert response.status_code == 400


def test_voice_preview_wraps_tts_error(tmp_path):
    def failing_preview(text, voice, rate=None, volume=None):
        raise TtsError("edge-tts streaming failed: network error")

    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        voice_preview_fn=failing_preview,
    )
    client = TestClient(app)

    response = client.get("/api/voice-preview/zh-TW-HsiaoChenNeural")
    assert response.status_code == 502


def test_voice_preview_forwards_rate_and_volume_and_caches_separately(tmp_path):
    calls = []

    def fake_preview(text, voice, rate=None, volume=None):
        calls.append((voice, rate, volume))
        return f"{rate}|{volume}".encode()

    manager = JobManager(pipeline_fn=lambda **kwargs: {}, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        voice_preview_fn=fake_preview,
    )
    client = TestClient(app)

    default = client.get("/api/voice-preview/zh-TW-HsiaoChenNeural")
    fast = client.get("/api/voice-preview/zh-TW-HsiaoChenNeural?rate=%2B20%25&volume=-10%25")
    fast_again = client.get("/api/voice-preview/zh-TW-HsiaoChenNeural?rate=%2B20%25&volume=-10%25")

    assert default.content == b"+0%|+0%"
    assert fast.content == b"+20%|-10%"
    assert fast_again.content == b"+20%|-10%"
    assert len(calls) == 2  # default and +20%/-10% each synthesized once, second +20% call cached


def test_gemini_api_key_reaches_pipeline_but_never_echoed_in_status(tmp_path):
    calls = []

    def fake_pipeline(**kwargs):
        calls.append(kwargs)
        return {}

    client, manager = _make_client(fake_pipeline, tmp_path)
    job_id = client.post(
        "/api/jobs",
        data=_upload_form(script_mode="AUTO", gemini_api_key="super-secret-key"),
        files=_upload_files(),
    ).json()["job_id"]
    manager.process_next()

    assert calls[0]["gemini_api_key"] == "super-secret-key"

    response = client.get(f"/api/jobs/{job_id}")
    assert "super-secret-key" not in response.text


# ---------- basic auth gate (opt-in, for exposing the app publicly) ----------


def _basic_auth_header(user, password):
    token = base64.b64encode(f"{user}:{password}".encode()).decode()
    return {"Authorization": f"Basic {token}"}


def test_no_auth_required_when_not_configured(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kw: None, auto_start=False)
    app = create_app(job_manager=manager, data_root=str(tmp_path / "data"), frontend_dist=None)
    client = TestClient(app)

    response = client.get("/api/jobs/does-not-exist")

    assert response.status_code == 404  # reached the route rather than being blocked by auth


def test_rejects_request_without_credentials_when_auth_configured(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kw: None, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        basic_auth_user="admin",
        basic_auth_password="s3cret",
    )
    client = TestClient(app)

    response = client.get("/api/jobs/does-not-exist")

    assert response.status_code == 401
    assert response.headers["www-authenticate"].startswith("Basic")


def test_rejects_wrong_credentials(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kw: None, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        basic_auth_user="admin",
        basic_auth_password="s3cret",
    )
    client = TestClient(app)

    response = client.get(
        "/api/jobs/does-not-exist", headers=_basic_auth_header("admin", "wrong-password")
    )

    assert response.status_code == 401


def test_accepts_correct_credentials(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kw: None, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        basic_auth_user="admin",
        basic_auth_password="s3cret",
    )
    client = TestClient(app)

    response = client.get(
        "/api/jobs/does-not-exist", headers=_basic_auth_header("admin", "s3cret")
    )

    assert response.status_code == 404  # past the auth gate, reached the route


def test_auth_gate_also_protects_the_mounted_frontend(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kw: None, auto_start=False)
    frontend_dir = tmp_path / "dist"
    frontend_dir.mkdir()
    (frontend_dir / "index.html").write_text("<html>app</html>", encoding="utf-8")
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=str(frontend_dir),
        basic_auth_user="admin",
        basic_auth_password="s3cret",
    )
    client = TestClient(app)

    unauthenticated = client.get("/")
    authenticated = client.get("/", headers=_basic_auth_header("admin", "s3cret"))

    assert unauthenticated.status_code == 401
    assert authenticated.status_code == 200


def test_download_route_is_exempt_from_basic_auth(tmp_path):
    # The mobile-download QR code points straight at this URL. A phone's
    # camera app / QR scanner typically opens it in a lightweight in-app
    # browser that doesn't render a Basic Auth login prompt for a direct
    # file download — the request just silently fails, which is what
    # actually reached us as "can scan the QR but can't download or view".
    # The job id itself is an unguessable UUID, so exempting only this one
    # route (not job creation or the rest of the app) keeps the QR
    # share-a-finished-video flow frictionless without weakening the gate
    # on the parts that cost real compute (uploading, generating).
    out_file = tmp_path / "課程.mp4"
    out_file.write_bytes(b"video-bytes")

    def fake_pipeline(**kwargs):
        return {"mp4": str(out_file)}

    manager = JobManager(pipeline_fn=fake_pipeline, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        basic_auth_user="admin",
        basic_auth_password="s3cret",
    )
    client = TestClient(app)
    job_id = client.post(
        "/api/jobs",
        data=_upload_form(),
        files=_upload_files(),
        headers=_basic_auth_header("admin", "s3cret"),
    ).json()["job_id"]
    manager.process_next()

    response = client.get(f"/api/jobs/{job_id}/download/mp4")  # no credentials

    assert response.status_code == 200
    assert response.content == b"video-bytes"


def test_job_status_and_job_creation_still_require_auth_when_configured(tmp_path):
    manager = JobManager(pipeline_fn=lambda **kw: None, auto_start=False)
    app = create_app(
        job_manager=manager,
        data_root=str(tmp_path / "data"),
        frontend_dist=None,
        basic_auth_user="admin",
        basic_auth_password="s3cret",
    )
    client = TestClient(app)

    assert client.get("/api/jobs/does-not-exist").status_code == 401
    assert client.post("/api/jobs", data=_upload_form(), files=_upload_files()).status_code == 401
